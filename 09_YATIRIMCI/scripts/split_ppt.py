from pptx import Presentation
import os
import os
BASE = os.path.dirname(__file__)
PPTX_PATH = os.path.join(BASE, "..", "GURBET_RADIO_PITCH_DECK.pptx")

def split_ppt(input_file, output_dir="../../BUILD_OUTPUT/slides"):
    prs = Presentation(input_file)
    os.makedirs(output_dir, exist_ok=True)

    for i, slide in enumerate(prs.slides):
        new_prs = Presentation()
        blank_slide_layout = new_prs.slide_layouts[6]
        new_slide = new_prs.slides.add_slide(blank_slide_layout)

        for shape in slide.shapes:
            el = shape.element
            new_slide.shapes._spTree.insert_element_before(el, 'p:extLst')

        new_prs.save(f"{output_dir}/slide_{i+1}.pptx")

if __name__ == "__main__":
    split_ppt(PPTX_PATH)

